import io
import re
from datetime import datetime

import numpy as np
import pandas as pd
import streamlit as st

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from docx import Document
except Exception:
    Document = None

from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

st.set_page_config(page_title="Win/Loss Intelligence Dashboard", page_icon="🏆", layout="wide")

SYNTHETIC_RFP = """
Synthetic LSH Managed IT Services RFP

The customer is seeking a long-term partner for Application Managed Services and Application Development across Commercial and Technical Operations. Scope includes Run services, selected Change mechanisms, transition, service integration, AIOps, DevOps, transformation, innovation, compliance, service levels, and commercial competitiveness.

The customer expects a compliant response, binding pricing, clear baseline service charges, transparent cost drivers, defensible delivery capacity, no financial engineering, and pricing that is not dependent on optional transformation initiatives. The customer may evaluate commercial attractiveness, credibility of baseline pricing, scope coverage, contract compliance, solution quality, transition confidence, technical capability, automation, service reliability, compliance, GxP understanding, staff retention, transformation and innovation commitments.

The customer expects vendors to avoid generic sales material and provide concrete solution detail, including technical solution overview, AIOps solution, delivery model, end-to-end ownership, service credit construct, niche resource capability, software ownership, DevOps delivery model, innovation plan, key assumptions, dependencies, transition plan, transformation plan, and termination assistance.
"""

SYNTHETIC_SOLUTION = """
Synthetic LSH Proposal Response

The vendor proposes an integrated and intelligent application services model with global delivery, AI-first operations, SRE-led service reliability, business observability, compliance at the core, transition governance, portfolio rationalization, and innovation. The solution includes Commercial and TechOps coverage across platforms such as SAP, Veeva, Salesforce, ServiceNow, Workday, cloud, data platforms, lab systems, manufacturing systems, QMS and integration layers.

The response emphasizes AI-driven predictive operations, AIOps, self-healing, ticket prevention, knowledge management, business observability, agentic L1.5, DevOps-integrated delivery, end-to-end ownership, transition waves, PMO governance, value realization office, transformation office, experience management office, innovation funds, and transformation funds.

The proposal is strong on solution themes, technology depth, AI/automation, domain narrative and transition confidence. However, the executive presentation may need tighter linkage between the solution architecture and the commercial model, clearer quantified baseline pricing logic, less generic sales messaging, and stronger executive-level explanation of how the solution changes cost, risk and outcomes.
"""

SYNTHETIC_CUSTOMER_FEEDBACK = """
Outcome: Lost

Customer feedback:
1. Gap between commercials and solution.
2. Executive session felt like more sales pitch and less solution.
3. Commercials were in the mid range.

Additional synthetic industry-style feedback:
Positive:
- Strong life sciences domain language and good understanding of regulated environments.
- Good coverage of AIOps, automation, observability and AI-enabled operations.
- Transition approach was structured and credible.
- Platform breadth across enterprise applications was appreciated.
- Governance model and service integration narrative were strong.

Negative:
- Commercial model was not sufficiently linked to solution levers, productivity commitments and baseline assumptions.
- Executive messaging was too broad and not enough anchored in the customer evaluation criteria.
- Pricing was acceptable but not compelling enough to create a clear downselect advantage.
- Value story did not clearly separate run savings, transformation savings, vendor investment and customer investment.
- Too many innovation claims were presented without enough proof of contractual commitment, timing and measurable impact.
- Some solution themes were strong but sounded generic compared with competitor-specific propositions.
- The response did not sufficiently show how commercial flexibility would work if scope was split across multiple vendors.
- Productivity improvements were directionally strong but needed clearer year-wise realization, accountability and risk-sharing.
- The customer expected more direct answers to how baseline pricing would remain defensible without relying on future transformation.
- Executive session needed a sharper problem-solution-commercial storyline, not a capabilities showcase.
"""

WIN_LOSS_DIMENSIONS = [
    {"id": "customer_priority_alignment", "name": "Customer Priority Alignment", "weight": 12, "positive": ["aligned", "customer priority", "evaluation criteria", "business objective", "sourcing objective", "specific", "direct answer", "requirement", "scope coverage"], "negative": ["misaligned", "generic", "not specific", "missed priority", "did not address", "unclear", "broad", "not anchored"], "recommendation": "Anchor every executive message and solution section to the customer's stated evaluation criteria, sourcing objectives and shortlisting logic."},
    {"id": "commercial_competitiveness", "name": "Commercial Competitiveness", "weight": 13, "positive": ["competitive", "pricing", "commercial", "baseline", "productivity", "risk sharing", "discount", "savings", "TCO", "defensible", "transparent"], "negative": ["mid range", "too expensive", "uncompetitive", "commercial gap", "pricing gap", "not compelling", "financial engineering", "not transparent"], "recommendation": "Create a clear commercial bridge from baseline cost to productivity, run savings, optional transformation, investments, risks and customer value."},
    {"id": "solution_commercial_linkage", "name": "Solution-to-Commercial Linkage", "weight": 12, "positive": ["solution linked", "commercial linkage", "savings lever", "MECE", "cost driver", "pricing structure", "productivity lever", "baseline adjustment"], "negative": ["gap between commercials and solution", "commercials and solution", "not linked", "disconnect", "unclear linkage", "solution gap"], "recommendation": "For every major solution lever, show the commercial implication: cost reduction, investment, timing, owner, risk and measurable benefit."},
    {"id": "executive_messaging", "name": "Executive Messaging Quality", "weight": 10, "positive": ["executive", "board", "outcome", "business value", "decision", "strategic", "clear story", "why us"], "negative": ["sales pitch", "too much sales", "marketing", "capabilities showcase", "not enough solution", "generic deck", "too broad"], "recommendation": "Use an executive storyline: customer problem, decision criteria, proposed answer, quantified value, commercial confidence and proof."},
    {"id": "differentiation_win_themes", "name": "Differentiation & Win Themes", "weight": 10, "positive": ["differentiator", "win theme", "unique", "competitive advantage", "why us", "proof point", "accelerator", "case study"], "negative": ["generic", "not differentiated", "similar to competitors", "unclear why", "no proof", "weak win theme"], "recommendation": "Convert capabilities into 3-5 client-specific win themes with proof, quantified benefit and competitor-aware positioning."},
    {"id": "ai_innovation_relevance", "name": "AI & Innovation Relevance", "weight": 9, "positive": ["AI", "GenAI", "AIOps", "automation", "self-healing", "observability", "agentic", "predictive", "innovation", "copilot"], "negative": ["AI unclear", "innovation vague", "not committed", "no roadmap", "vendor lock-in", "unclear ownership", "hype"], "recommendation": "Tie AI use cases to specific operational outcomes, adoption roadmap, data/security model, contractual commitments and measurable productivity."},
    {"id": "domain_credibility", "name": "LSH Domain Credibility", "weight": 8, "positive": ["GxP", "CSV", "CSA", "validation", "pharma", "life sciences", "clinical", "regulatory", "manufacturing", "quality", "patient"], "negative": ["domain gap", "insufficient domain", "not enough pharma", "weak regulatory", "weak validation", "no LSH proof"], "recommendation": "Strengthen LSH proof points by function: Commercial, TechOps, Quality, Manufacturing, Clinical, Regulatory, Safety and Patient."},
    {"id": "delivery_confidence", "name": "Delivery Confidence", "weight": 8, "positive": ["transition", "governance", "PMO", "ramp-up", "resource", "site", "wave", "risk mitigation", "RACI", "SLA"], "negative": ["delivery risk", "unclear transition", "weak governance", "resourcing concern", "ramp-up risk", "unclear accountability"], "recommendation": "Show day-0 to steady-state execution with named governance, resourcing, risks, mitigation, dependencies, and measurable readiness gates."},
    {"id": "proof_references", "name": "Proof Points & References", "weight": 7, "positive": ["case study", "reference", "benchmark", "proven", "track record", "comparable", "experience", "similar client"], "negative": ["no proof", "unsupported", "claim", "not evidenced", "insufficient reference", "weak case study"], "recommendation": "Back major claims with relevant life sciences examples, benchmarks, before/after metrics and named reusable assets."},
    {"id": "rfp_responsiveness", "name": "RFP Responsiveness & Compliance", "weight": 6, "positive": ["compliant", "response", "format", "instructions", "mandatory", "binding", "contract", "assumptions", "dependencies"], "negative": ["non-compliant", "missing", "not submitted", "wrong format", "deferred", "assumption", "change control"], "recommendation": "Track every RFP instruction and evaluation criterion to a response artifact, owner and evidence location."},
    {"id": "future_repeatability", "name": "Repeatability for Future Deals", "weight": 5, "positive": ["repeat", "playbook", "replicable", "reuse", "template", "standard", "pattern", "learning"], "negative": ["one-off", "not reusable", "no learning", "not captured", "lost lesson"], "recommendation": "Convert this deal's learnings into a reusable pursuit playbook: solution, commercial, executive messaging and objection handling."},
]


def clean_text(text):
    return re.sub(r"\s+", " ", text or "").strip()


def extract_pdf(file_obj):
    if PdfReader is None:
        raise RuntimeError("pypdf is not installed. Add pypdf to requirements.txt.")
    reader = PdfReader(file_obj)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_pptx(file_obj):
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Add python-pptx to requirements.txt.")
    prs = Presentation(file_obj)
    chunks = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks.append(f"\n--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    chunks.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(chunks)


def extract_docx(file_obj):
    if Document is None:
        raise RuntimeError("python-docx is not installed. Add python-docx to requirements.txt.")
    doc = Document(file_obj)
    chunks = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            chunks.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(chunks)


def extract_text(uploaded_file):
    name = uploaded_file.name.lower()
    data = uploaded_file.read()
    buffer = io.BytesIO(data)
    if name.endswith(".pdf"):
        return extract_pdf(buffer)
    if name.endswith(".pptx"):
        return extract_pptx(buffer)
    if name.endswith(".docx"):
        return extract_docx(buffer)
    if name.endswith(".txt"):
        return data.decode("utf-8", errors="ignore")
    raise ValueError("Unsupported file type. Upload PDF, PPTX, DOCX, or TXT.")


def keyword_hits(text, keywords):
    text_lower = text.lower()
    found = []
    for keyword in keywords:
        if keyword.lower() in text_lower:
            found.append(keyword)
    return len(found), found


def cosine_score(query_terms, text):
    query = " ".join(query_terms)
    docs = [query, text]
    try:
        vectorizer = TfidfVectorizer(stop_words="english", ngram_range=(1, 2), max_features=12000)
        matrix = vectorizer.fit_transform(docs)
        return float(cosine_similarity(matrix[0:1], matrix[1:2])[0][0])
    except Exception:
        return 0.0


def score_win_loss_dimension(rfp_text, solution_text, feedback_text, dimension):
    positive_terms = dimension["positive"]
    negative_terms = dimension["negative"]
    combined_solution = solution_text + " " + rfp_text
    pos_count, pos_hits = keyword_hits(combined_solution, positive_terms)
    neg_count, neg_hits = keyword_hits(feedback_text, negative_terms)
    pos_sem = cosine_score(positive_terms, combined_solution)
    neg_sem = cosine_score(negative_terms, feedback_text)
    strength = min(1.0, (pos_count / max(3, min(8, len(positive_terms)))) * 0.7 + min(1.0, pos_sem * 6) * 0.3)
    friction = min(1.0, (neg_count / max(2, min(5, len(negative_terms)))) * 0.7 + min(1.0, neg_sem * 6) * 0.3)
    raw = (0.65 * strength + 0.35 * (1 - friction)) * 5
    score_0_to_5 = round(max(0.0, min(5.0, raw)), 2)
    weighted = round((score_0_to_5 / 5.0) * dimension["weight"], 2)
    if friction >= 0.65:
        impact = "Likely loss driver"
    elif friction >= 0.35:
        impact = "Possible improvement area"
    elif strength >= 0.65:
        impact = "Likely win/supporting driver"
    else:
        impact = "Low evidence / needs review"
    return {"Dimension": dimension["name"], "Weight": dimension["weight"], "Score / 5": score_0_to_5, "Weighted Score": weighted, "Post-Facto Impact": impact, "Positive Signals": ", ".join(pos_hits[:10]) if pos_hits else "None detected", "Negative Feedback Signals": ", ".join(neg_hits[:10]) if neg_hits else "None detected", "Recommended Next-Deal Action": dimension["recommendation"]}


def analyze_feedback_themes(feedback_text):
    positive_markers = ["strong", "appreciated", "good", "credible", "clear", "well", "positive", "liked", "confidence", "differentiated", "excellent", "compelling"]
    negative_markers = ["gap", "weak", "less", "mid range", "generic", "unclear", "not enough", "concern", "risk", "lost", "eliminated", "uncompetitive", "disconnect", "sales pitch", "not compelling"]
    _, positives = keyword_hits(feedback_text, positive_markers)
    _, negatives = keyword_hits(feedback_text, negative_markers)
    return positives, negatives


def score_deal(rfp_text, solution_text, feedback_text):
    rows = [score_win_loss_dimension(rfp_text, solution_text, feedback_text, dim) for dim in WIN_LOSS_DIMENSIONS]
    scorecard = pd.DataFrame(rows)
    total = round(float(scorecard["Weighted Score"].sum()), 2)
    likely_loss_drivers = scorecard[scorecard["Post-Facto Impact"].isin(["Likely loss driver", "Possible improvement area"])].sort_values(["Score / 5", "Weight"], ascending=[True, False]).head(5)
    likely_win_drivers = scorecard[scorecard["Post-Facto Impact"].isin(["Likely win/supporting driver"])].sort_values(["Score / 5", "Weight"], ascending=[False, False]).head(5)
    return scorecard, total, likely_win_drivers, likely_loss_drivers


def build_executive_summary(outcome, total, loss_drivers, win_drivers):
    if outcome == "Won":
        lead = "The pursuit appears to have converted because the response created enough customer confidence across the highest-value decision themes."
        action = "Codify the winning pattern into the next-deal playbook and reuse the strongest win themes, proof points and commercial levers."
    elif outcome == "Lost":
        lead = "The pursuit appears to have been eliminated due to gaps between customer decision priorities, commercial attractiveness and executive-level solution clarity."
        action = "Prioritize corrective actions before the next pursuit: sharpen commercial linkage, reduce sales-pitch language, strengthen quantified outcomes and make the executive story more solution-led."
    elif outcome == "No Decision":
        lead = "The pursuit did not convert into a decision, which may indicate unclear value, timing, business case, sponsorship, or commercial confidence."
        action = "Strengthen urgency, business case, outcome quantification and executive sponsorship in future pursuits."
    else:
        lead = "The pursuit requires additional customer feedback before a high-confidence win/loss root cause can be assigned."
        action = "Capture structured debrief notes and rerun the dashboard with updated feedback."
    top_loss = ", ".join(loss_drivers["Dimension"].head(3).tolist()) if len(loss_drivers) else "No clear loss drivers detected"
    top_win = ", ".join(win_drivers["Dimension"].head(3).tolist()) if len(win_drivers) else "No clear win drivers detected"
    return lead, top_loss, top_win, action


def create_history_record(deal_name, client, outcome, total, loss_drivers):
    return {"Date Added": datetime.now().strftime("%Y-%m-%d %H:%M"), "Deal Name": deal_name, "Client": client, "Outcome": outcome, "Win/Loss Intelligence Score": total, "Top Driver 1": loss_drivers.iloc[0]["Dimension"] if len(loss_drivers) > 0 else "", "Top Driver 2": loss_drivers.iloc[1]["Dimension"] if len(loss_drivers) > 1 else "", "Top Driver 3": loss_drivers.iloc[2]["Dimension"] if len(loss_drivers) > 2 else ""}


st.title("🏆 Win/Loss Intelligence & Pursuit Improvement Dashboard")
st.caption("Post-facto root-cause analysis across RFP, submitted solution and customer feedback. Designed for LSH managed services, ASM, AD, modernization, AI, data, SAP, Salesforce, Veeva, Workday and related pursuits.")

with st.sidebar:
    st.header("Operating Mode")
    use_synthetic = st.toggle("Use synthetic demo data", value=True)
    st.header("Deal Context")
    deal_name = st.text_input("Deal name", value="Roche Cluster 2 Managed IT Services")
    client_name = st.text_input("Client / account", value="LSH Client")
    outcome = st.selectbox("Deal outcome", ["Lost", "Won", "No Decision", "Down-selected", "Pending Debrief"], index=0)
    competitor = st.text_input("Known winner / competitor", value="")
    deal_type = st.multiselect("Deal scope", ["ASM", "AD", "Application Modernization", "AIOps", "Data & Analytics", "SAP", "Salesforce", "Veeva", "Workday", "ServiceNow", "LSH Managed Services"], default=["ASM", "AD", "AIOps", "LSH Managed Services"])
    st.divider()
    st.header("Win/Loss Model")
    st.write(f"Total weight: **{sum(d['weight'] for d in WIN_LOSS_DIMENSIONS)}**")
    st.dataframe(pd.DataFrame([{"Dimension": d["name"], "Weight": d["weight"]} for d in WIN_LOSS_DIMENSIONS]), use_container_width=True, hide_index=True)
    st.divider()
    st.header("Historical Learning")
    history_file = st.file_uploader("Optional: upload prior win/loss history CSV", type=["csv"])

st.subheader("1. Input Artifacts")
if use_synthetic:
    rfp_text = clean_text(SYNTHETIC_RFP)
    solution_text = clean_text(SYNTHETIC_SOLUTION)
    feedback_text = clean_text(SYNTHETIC_CUSTOMER_FEEDBACK)
    st.success("Using embedded synthetic RFP, solution and customer feedback.")
else:
    c1, c2, c3 = st.columns(3)
    with c1:
        rfp_file = st.file_uploader("Upload RFP", type=["pdf", "pptx", "docx", "txt"])
    with c2:
        solution_file = st.file_uploader("Upload submitted solution / proposal", type=["pdf", "pptx", "docx", "txt"])
    with c3:
        feedback_file = st.file_uploader("Optional: upload feedback notes", type=["pdf", "pptx", "docx", "txt"])
    feedback_manual = st.text_area("Paste customer debrief / evaluator feedback", value="Gap between commercials and solution.\nExec session - we did more sales pitch, less solution.\nYour commercials are in mid range.", height=140)
    if not rfp_file or not solution_file:
        st.warning("Upload both the RFP and submitted solution, or turn synthetic mode ON.")
        st.stop()
    with st.spinner("Parsing uploaded artifacts..."):
        rfp_text = clean_text(extract_text(rfp_file))
        solution_text = clean_text(extract_text(solution_file))
        feedback_parts = [feedback_manual]
        if feedback_file is not None:
            feedback_parts.append(extract_text(feedback_file))
        feedback_text = clean_text("\n".join(feedback_parts))

st.subheader("2. Executive Win/Loss Readout")
scorecard, total, win_drivers, loss_drivers = score_deal(rfp_text, solution_text, feedback_text)
lead, top_loss, top_win, action = build_executive_summary(outcome, total, loss_drivers, win_drivers)

m1, m2, m3, m4 = st.columns(4)
m1.metric("Win/Loss Intelligence Score", f"{total} / 100")
m2.metric("Outcome", outcome)
m3.metric("Likely Loss Drivers", len(loss_drivers))
m4.metric("Likely Win Drivers", len(win_drivers))

if outcome == "Won":
    st.success(lead)
elif outcome == "Lost":
    st.error(lead)
else:
    st.warning(lead)

st.markdown(f"**Top loss / improvement drivers:** {top_loss}")
st.markdown(f"**Top win / supporting drivers:** {top_win}")
st.info(f"**Next-deal executive action:** {action}")

st.subheader("3. Root-Cause Scorecard")
st.dataframe(scorecard, use_container_width=True, hide_index=True)

st.subheader("4. Win/Loss Drivers")
col_a, col_b = st.columns(2)
with col_a:
    st.markdown("### Likely Loss / Improvement Drivers")
    if len(loss_drivers):
        for _, row in loss_drivers.iterrows():
            st.markdown(f"**{row['Dimension']}** — `{row['Score / 5']} / 5`")
            st.write(row["Recommended Next-Deal Action"])
    else:
        st.write("No clear loss driver detected from the current feedback.")
with col_b:
    st.markdown("### Likely Win / Supporting Drivers")
    if len(win_drivers):
        for _, row in win_drivers.iterrows():
            st.markdown(f"**{row['Dimension']}** — `{row['Score / 5']} / 5`")
            st.write("Reuse this theme in future pursuits with stronger evidence and customer-specific framing.")
    else:
        st.write("No clear win/supporting driver detected from the current feedback.")

st.subheader("5. Customer Feedback Theme Scan")
pos_signals, neg_signals = analyze_feedback_themes(feedback_text)
c1, c2 = st.columns(2)
with c1:
    st.markdown("### Positive Feedback Signals")
    st.write(", ".join(pos_signals) if pos_signals else "No explicit positive sentiment markers detected.")
with c2:
    st.markdown("### Negative Feedback Signals")
    st.write(", ".join(neg_signals) if neg_signals else "No explicit negative sentiment markers detected.")

st.subheader("6. Pursuit Improvement Playbook")
playbook_rows = []
for _, row in loss_drivers.iterrows():
    playbook_rows.append({"Priority": len(playbook_rows) + 1, "Area": row["Dimension"], "Action": row["Recommended Next-Deal Action"], "Owner": "Solution / Sales / Commercial / Domain SME", "When to Apply": "Next similar LSH pursuit"})
if playbook_rows:
    st.dataframe(pd.DataFrame(playbook_rows), use_container_width=True, hide_index=True)
else:
    st.write("No prioritized corrective actions available yet.")

st.subheader("7. Evolving Knowledge Base")
history_df = pd.DataFrame()
if history_file is not None:
    try:
        history_df = pd.read_csv(history_file)
    except Exception as e:
        st.warning(f"Could not read uploaded history CSV: {e}")
current_record = create_history_record(deal_name, client_name, outcome, total, loss_drivers)
updated_history = pd.concat([history_df, pd.DataFrame([current_record])], ignore_index=True)
st.write("Each new deal can be added to the learning base by downloading this CSV and uploading it in the next run.")
st.dataframe(updated_history, use_container_width=True, hide_index=True)
st.download_button("Download updated win/loss history CSV", data=updated_history.to_csv(index=False).encode("utf-8"), file_name="win_loss_learning_history.csv", mime="text/csv")

st.subheader("8. Export Executive Scorecard")
output = io.BytesIO()
with pd.ExcelWriter(output, engine="openpyxl") as writer:
    pd.DataFrame([{"Deal Name": deal_name, "Client": client_name, "Outcome": outcome, "Competitor": competitor, "Deal Scope": ", ".join(deal_type), "Score": total, "Executive Summary": lead, "Top Loss Drivers": top_loss, "Top Win Drivers": top_win, "Next Action": action}]).to_excel(writer, index=False, sheet_name="Executive Summary")
    scorecard.to_excel(writer, index=False, sheet_name="Root Cause Scorecard")
    loss_drivers.to_excel(writer, index=False, sheet_name="Loss Drivers")
    win_drivers.to_excel(writer, index=False, sheet_name="Win Drivers")
    pd.DataFrame(playbook_rows).to_excel(writer, index=False, sheet_name="Improvement Playbook")
    updated_history.to_excel(writer, index=False, sheet_name="Learning History")
st.download_button("Download Excel Win/Loss Analysis", data=output.getvalue(), file_name="win_loss_intelligence_analysis.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

st.subheader("9. Artifact Preview")
with st.expander("RFP text preview"):
    st.write(rfp_text[:8000])
with st.expander("Solution text preview"):
    st.write(solution_text[:8000])
with st.expander("Customer feedback preview"):
    st.write(feedback_text[:8000])

